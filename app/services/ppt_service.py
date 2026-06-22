"""Build a .pptx from an Outline using python-pptx.

Style templates control colors, fonts, and decorative shapes per `style`.
Layout values map to a small set of slide builders.
"""
from __future__ import annotations

from pathlib import Path
from typing import Callable

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

from app.models import Outline
from app.services.image_service import ImageService


STYLE_TEMPLATES: dict[str, dict] = {
    "business": {
        "title_color": RGBColor(0x1F, 0x3A, 0x68),
        "body_color": RGBColor(0x33, 0x33, 0x33),
        "accent_color": RGBColor(0x2E, 0x86, 0xC1),
        "title_font": "Microsoft YaHei",
        "body_font": "Microsoft YaHei",
        "bg_color": RGBColor(0xFF, 0xFF, 0xFF),
    },
    "academic": {
        "title_color": RGBColor(0x22, 0x22, 0x22),
        "body_color": RGBColor(0x33, 0x33, 0x33),
        "accent_color": RGBColor(0x8B, 0x6F, 0x4E),
        "title_font": "SimSun",
        "body_font": "SimSun",
        "bg_color": RGBColor(0xFA, 0xF6, 0xEE),
    },
    "minimal": {
        "title_color": RGBColor(0x00, 0x00, 0x00),
        "body_color": RGBColor(0x44, 0x44, 0x44),
        "accent_color": RGBColor(0x00, 0x00, 0x00),
        "title_font": "Helvetica",
        "body_font": "Helvetica",
        "bg_color": RGBColor(0xFF, 0xFF, 0xFF),
    },
    "creative": {
        "title_color": RGBColor(0xE6, 0x4A, 0xC0),
        "body_color": RGBColor(0x33, 0x33, 0x33),
        "accent_color": RGBColor(0x4C, 0xC9, 0xF0),
        "title_font": "Source Han Sans CN",
        "body_font": "Source Han Sans CN",
        "bg_color": RGBColor(0xFD, 0xF6, 0xFF),
    },
}


class PPTService:
    def __init__(self, image_service: ImageService | None = None) -> None:
        self._images = image_service or ImageService()

    def build(
        self,
        *,
        outline: Outline,
        style: str,
        image_mode: str,
        out_path: Path,
        images_dir: Path,
    ) -> Path:
        template = STYLE_TEMPLATES.get(style, STYLE_TEMPLATES["business"])
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        images_dir.mkdir(parents=True, exist_ok=True)

        builders: dict[str, Callable] = {
            "title": self._build_title_slide,
            "title-content": self._build_title_content,
            "two-column": self._build_two_column,
            "quote": self._build_quote,
            "section": self._build_section,
        }
        for idx, page in enumerate(outline.pages):
            builder = builders.get(page.layout, self._build_title_content)
            builder(prs, outline, page, idx, template, image_mode, images_dir)
        out_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(out_path)
        return out_path

    # -- builders ---------------------------------------------------------
    def _add_blank(self, prs):
        blank = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank)
        slide.background.fill.solid()
        return slide

    def _build_title_slide(self, prs, outline, page, idx, t, image_mode, images_dir):
        slide = self._add_blank(prs)
        slide.background.fill.fore_color.rgb = t["bg_color"]
        tx = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(2))
        tf = tx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = outline.topic
        r.font.size = Pt(54)
        r.font.bold = True
        r.font.color.rgb = t["title_color"]
        r.font.name = t["title_font"]

    def _build_title_content(self, prs, outline, page, idx, t, image_mode, images_dir):
        slide = self._add_blank(prs)
        slide.background.fill.fore_color.rgb = t["bg_color"]
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(1))
        title_tf = title_box.text_frame
        title_tf.text = page.title
        title_tf.paragraphs[0].runs[0].font.size = Pt(36)
        title_tf.paragraphs[0].runs[0].font.bold = True
        title_tf.paragraphs[0].runs[0].font.color.rgb = t["title_color"]
        title_tf.paragraphs[0].runs[0].font.name = t["title_font"]

        body_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.5), Inches(5))
        body_tf = body_box.text_frame
        body_tf.word_wrap = True
        for i, point in enumerate(page.key_points):
            p = body_tf.paragraphs[0] if i == 0 else body_tf.add_paragraph()
            p.text = f"• {point}"
            p.runs[0].font.size = Pt(22)
            p.runs[0].font.color.rgb = t["body_color"]
            p.runs[0].font.name = t["body_font"]
            p.space_after = Pt(8)

        if image_mode == "placeholder":
            self._images.placeholder(
                out_path=images_dir / f"slide_{idx}.png",
                color=(245, 245, 245),
                text=f"图 {idx + 1}",
            )
            slide.shapes.add_picture(
                str(images_dir / f"slide_{idx}.png"),
                Inches(9.5), Inches(5.0), Inches(3.3), Inches(1.8),
            )
        elif image_mode == "ai":
            try:
                self._images.ai_generate(
                    out_path=images_dir / f"slide_{idx}.png",
                    prompt=f"Illustration for slide titled '{page.title}', {outline.topic} topic",
                )
                slide.shapes.add_picture(
                    str(images_dir / f"slide_{idx}.png"),
                    Inches(9.5), Inches(5.0), Inches(3.3), Inches(1.8),
                )
            except RuntimeError:
                # Image provider not configured or call failed — fall back silently.
                pass

    def _build_two_column(self, prs, outline, page, idx, t, image_mode, images_dir):
        slide = self._add_blank(prs)
        slide.background.fill.fore_color.rgb = t["bg_color"]
        title = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(1))
        title.text_frame.text = page.title
        title.text_frame.paragraphs[0].runs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].runs[0].font.bold = True
        title.text_frame.paragraphs[0].runs[0].font.color.rgb = t["title_color"]
        title.text_frame.paragraphs[0].runs[0].font.name = t["title_font"]
        left = page.key_points[: len(page.key_points) // 2 or 1]
        right = page.key_points[len(left):]
        for col_x, items in ((Inches(0.6), left), (Inches(7.0), right)):
            box = slide.shapes.add_textbox(col_x, Inches(1.8), Inches(5.8), Inches(5))
            tf = box.text_frame
            tf.word_wrap = True
            for i, point in enumerate(items or [""]):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = f"• {point}"
                p.runs[0].font.size = Pt(22)
                p.runs[0].font.color.rgb = t["body_color"]
                p.runs[0].font.name = t["body_font"]
                p.space_after = Pt(8)

    def _build_quote(self, prs, outline, page, idx, t, image_mode, images_dir):
        slide = self._add_blank(prs)
        slide.background.fill.fore_color.rgb = t["bg_color"]
        text = page.key_points[0] if page.key_points else page.title
        quote = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(2.5))
        tf = quote.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"“{text}”"
        p.alignment = 2  # center
        p.runs[0].font.size = Pt(40)
        p.runs[0].font.italic = True
        p.runs[0].font.color.rgb = t["accent_color"]
        p.runs[0].font.name = t["body_font"]
        sub = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11.333), Inches(0.6))
        sub.text_frame.text = page.title
        sub.text_frame.paragraphs[0].runs[0].font.size = Pt(20)
        sub.text_frame.paragraphs[0].runs[0].font.color.rgb = t["body_color"]
        sub.text_frame.paragraphs[0].runs[0].font.name = t["body_font"]

    def _build_section(self, prs, outline, page, idx, t, image_mode, images_dir):
        slide = self._add_blank(prs)
        slide.background.fill.fore_color.rgb = t["title_color"]
        box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11.333), Inches(1.5))
        box.text_frame.text = page.title
        box.text_frame.paragraphs[0].alignment = 2
        box.text_frame.paragraphs[0].runs[0].font.size = Pt(48)
        box.text_frame.paragraphs[0].runs[0].font.bold = True
        box.text_frame.paragraphs[0].runs[0].font.color.rgb = t["bg_color"]
        box.text_frame.paragraphs[0].runs[0].font.name = t["title_font"]
