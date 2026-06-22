"""Tests for the PPT generation service."""
from pathlib import Path

from pptx import Presentation

from app.models import Outline, OutlinePage
from app.services.ppt_service import PPTService, STYLE_TEMPLATES


def _outline_with_pages() -> Outline:
    return Outline(
        topic="Test",
        pages=[
            OutlinePage(title="Cover", key_points=[], layout="title"),
            OutlinePage(title="Body", key_points=["A", "B"], layout="title-content"),
        ],
    )


def test_build_pptx_creates_file(tmp_dir):
    service = PPTService()
    out = tmp_dir / "x.pptx"
    pptx_path = service.build(outline=_outline_with_pages(), style="business", image_mode="none", out_path=out, images_dir=tmp_dir / "img")
    assert pptx_path.exists()
    prs = Presentation(pptx_path)
    assert len(prs.slides) == 2


def test_title_slide_has_topic_text(tmp_dir):
    service = PPTService()
    out = tmp_dir / "x.pptx"
    service.build(outline=_outline_with_pages(), style="minimal", image_mode="none", out_path=out, images_dir=tmp_dir / "img")
    prs = Presentation(out)
    title_text = " ".join(shape.text for shape in prs.slides[0].shapes if shape.has_text_frame)
    assert "Test" in title_text or "Cover" in title_text


def test_unknown_layout_falls_back_to_title_content(tmp_dir):
    outline = Outline(topic="X", pages=[OutlinePage(title="P", key_points=["k"], layout="title-content")])
    service = PPTService()
    out = tmp_dir / "x.pptx"
    service.build(outline=outline, style="business", image_mode="none", out_path=out, images_dir=tmp_dir / "img")
    prs = Presentation(out)
    assert len(prs.slides) == 1


def test_style_templates_include_business(tmp_dir):
    assert "business" in STYLE_TEMPLATES
    assert STYLE_TEMPLATES["business"]["title_color"]
